from pptx.slide import Slides, Slide
from .powerpoint_generator import PowerpointGenerator
from .util import _find_sunday, _find_slide_with_name, _find_name_in_iterable, _remove_element
from pptx.shapes.shapetree import GroupShapes, Shape
from datetime import datetime, timedelta
from pptx.dml.fill import FillFormat
from pptx.dml.color import RGBColor


def generate_eight_week(pg: PowerpointGenerator, slides: Slides):
    start_date = _find_sunday(pg.config.get_date())

    eight_week_slide: Slide = _find_slide_with_name(slides, 'eight-week')

    for i in range(0, 8):
        week_group = _find_name_in_iterable(eight_week_slide.shapes, 'week' + str(i+1))
        process_week(pg, week_group.shapes, start_date + timedelta(weeks=i), (i == 0))


def process_week(pg: PowerpointGenerator, group: GroupShapes, start_date: datetime, first_week: bool):
    show_month_name = first_week
    for i in range(0, 7):
        date = start_date + timedelta(days=i)
        day_block: Shape = _find_name_in_iterable(group, 'day' + str(i+1))
        day_block.text_frame.paragraphs[0].runs[0].text = str((start_date + timedelta(days=i)).day)
        day_fill: FillFormat = day_block.fill
        if date.month % 2 == 0:
            day_fill.fore_color.rgb = RGBColor.from_string(pg.config.eight_week.month_one_color)
        else:
            day_fill.fore_color.rgb = RGBColor.from_string(pg.config.eight_week.month_two_color)
        if date.day == 1:
            show_month_name = True
    month_block: Shape = _find_name_in_iterable(group, 'month')
    if show_month_name:
        month_block.text_frame.paragraphs[0].runs[0].text = (start_date + timedelta(days=6)).strftime('%B')
    else:
        _remove_element(month_block)
