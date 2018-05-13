from .powerpoint_generator import PowerpointGenerator
from .two_week_generator import generate_two_week
from .eight_week_generator import generate_eight_week
from pptx import Presentation


def generate_powerpoint(pg: PowerpointGenerator):
    pptx = Presentation(pg.config.get_template_path())

    #two_week_slide = pptx.slides[2]
    generate_two_week(pg, pptx.slides)

    generate_eight_week(pg, pptx.slides)

    pptx.save(pg.config.get_save_path())
