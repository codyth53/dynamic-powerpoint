from pptx.slide import Slide, Slides
from pptx.shapes.shapetree import GroupShapes, Shape, GroupShape
from .powerpoint_generator import PowerpointGenerator
from ..calendar_utils import CalendarEvent
from datetime import datetime, timedelta
from typing import List
from .util import _hide_slide, _find_slide_with_name, _find_name_in_iterable, _find_sunday
import logging
from pptx.dml.color import RGBColor


logger = logging.getLogger('dynamic_powerpoint.powerpoint_generator.two_week_generator')


def generate_two_week(pg: PowerpointGenerator, slides: Slides):
    start_date = _find_sunday(pg.config.get_date())

    events = pg.calendar.get_events(start=start_date, end=start_date + timedelta(days=14))
    unique_events = create_array_of_unique_events(events)

    groups_of_three = []
    while len(unique_events) > 0:
        if len(unique_events) > 3:
            group = unique_events[:3]
            groups_of_three.append(group)
            unique_events = unique_events[3:]
        else:
            groups_of_three.append(unique_events)
            unique_events = []

    logger.info('Found the following events:')
    for group in groups_of_three:
        logger.info('  Group:')
        for event in group:
            logger.info('    {0}'.format(str(event)))

    for i in range(1, 4):
        slide = _find_slide_with_name(slides, "two-week-" + str(i))
        if slide:
            if len(groups_of_three) >= i:
                generate_two_week_slide(slide, events, groups_of_three[i-1], start_date)
            else:
                _hide_slide(slide)


def generate_two_week_slide(slide: Slide,
                            all_events: List[CalendarEvent],
                            event_group: List[CalendarEvent],
                            start_date: datetime):
    two_week_calendar_group = _find_name_in_iterable(slide.shapes, 'dates')
    generate_calendar_dates(two_week_calendar_group.shapes, start_date)

    event_marker_group = _find_name_in_iterable(slide.shapes, 'event-markers')
    generate_event_markers(event_marker_group.shapes, all_events, start_date)

    event_blocks_group = _find_name_in_iterable(slide.shapes, 'event-blocks')
    generate_event_detail_boxes(event_blocks_group.shapes, event_group)


def generate_calendar_dates(group: GroupShapes, start_date: datetime):
    date = start_date
    textbox: Shape = None
    for textbox in group:
        # assuming textboxes with name 'dayX'
        num = int(textbox.name[3:])
        this_day = date + timedelta(days=(num-1))
        textbox.text_frame.paragraphs[0].runs[0].text = str(this_day.day)


def generate_event_markers(group: GroupShapes, events: List[CalendarEvent], start_date: datetime):
    for event_group in group:
        # assuming group with name 'eventX'
        num = int(event_group.name[5:])
        this_day = start_date + timedelta(days=(num-1))
        this_day_events = [x for x in events if x.start >= this_day and x.end <= this_day + timedelta(days=1)]
        if len(this_day_events) == 0:
            logger.debug('Day {0} has {0} events'.format(this_day, len(this_day_events)))
        else:
            logger.info('Day {0} has {1} events'.format(this_day, len(this_day_events)))
        if len(this_day_events) < 2:
            second_marker = _find_name_in_iterable(event_group.shapes, 'event2')
            ele = second_marker.element
            ele.getparent().remove(ele)
        if len(this_day_events) < 1:
            first_marker = _find_name_in_iterable(event_group.shapes, 'event1')
            ele = first_marker.element
            ele.getparent().remove(ele)


def generate_event_detail_boxes(group: GroupShapes, events: List[CalendarEvent]):
    for i in range(0, 3):
        event_block: GroupShape = _find_name_in_iterable(group, 'event-' + str(i+1))
        if len(events) > i:
            # handle block color
            textbox = _find_name_in_iterable(event_block.shapes, 'text')
            textbox.text = events[i].name

            if events[i].category is not None:
                event_shape: Shape = _find_name_in_iterable(event_block.shapes, 'shape')
                category = events[i].category
                fill = event_shape.fill
                fill.fore_color.rgb = RGBColor.from_string(category.color)
                if category.transparency is not None:
                    fill.transparency = category.transparency
        else:
            ele = event_block.element
            ele.getparent().remove(ele)


def create_array_of_unique_events(events: List[CalendarEvent]) -> List[CalendarEvent]:
    events_list = []
    recurring_ids = []

    for event in events:
        if event.repeat_tag:
            if event.repeat_tag in recurring_ids:
                continue
            recurring_ids.append(event.repeat_tag)
        events_list.append(event)

    return events_list
